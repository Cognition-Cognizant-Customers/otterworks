"""Realistic file-content generators for the RetailCo enterprise drive.

Each function returns ``(bytes, mime_type)`` for a given logical file so the
uploaded objects are real, openable files of the correct type rather than empty
placeholders. Content is derived deterministically from the file name + a seeded
``random.Random`` so re-runs are reproducible.

Heavy office formats (xlsx/docx/pptx/pdf/png/jpg) use optional third-party
libraries. If a library is missing the generator degrades gracefully to a
plain-text stand-in with the correct extension so the drive still populates.
"""
from __future__ import annotations

import base64
import io
import json
import random
from datetime import datetime, timedelta

# ---- optional heavy deps (degrade gracefully) -------------------------------
try:
    from openpyxl import Workbook
except Exception:  # pragma: no cover
    Workbook = None
try:
    import docx
except Exception:  # pragma: no cover
    docx = None
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None
try:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.pdfgen import canvas as _pdfcanvas
except Exception:  # pragma: no cover
    _pdfcanvas = None
try:
    from PIL import Image, ImageDraw
except Exception:  # pragma: no cover
    Image = None

MIME = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "csv": "text/csv",
    "txt": "text/plain",
    "md": "text/markdown",
    "json": "application/json",
    "png": "image/png",
    "jpg": "image/jpeg",
    "mp4": "video/mp4",
    "mp3": "audio/mpeg",
}

_PRODUCTS = [
    "Aurora Wireless Earbuds", "Nomad Trail Backpack", "Cirrus Down Jacket",
    "Harbor Cast-Iron Skillet", "Vertex Running Shoe", "Lumen LED Desk Lamp",
    "Terra Ceramic Mug Set", "Pulse Fitness Tracker", "Drift Cotton Bedsheets",
    "Summit Insulated Bottle", "Grove Bamboo Cutting Board", "Echo Bluetooth Speaker",
]
_REGIONS = ["Northeast", "Southeast", "Midwest", "Southwest", "West", "Pacific Northwest"]
_STORES = [f"Store #{1000 + i}" for i in range(60)]


def _rng(name: str, seed: int) -> random.Random:
    return random.Random(f"{seed}:{name}")


def _money(r: random.Random, lo: float, hi: float) -> float:
    return round(r.uniform(lo, hi), 2)


# ---- individual format builders --------------------------------------------
def _xlsx(name: str, r: random.Random) -> bytes:
    if Workbook is None:
        return _txt_fallback(name, r)
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    headers = ["Date", "Region", "Store", "Product", "Units", "Revenue", "Margin %"]
    ws.append(headers)
    start = datetime(2025, 1, 1)
    for i in range(r.randint(40, 200)):
        d = start + timedelta(days=r.randint(0, 480))
        units = r.randint(1, 800)
        rev = round(units * _money(r, 4.0, 240.0), 2)
        ws.append([
            d.strftime("%Y-%m-%d"), r.choice(_REGIONS), r.choice(_STORES),
            r.choice(_PRODUCTS), units, rev, round(r.uniform(8, 62), 1),
        ])
    # a small summary sheet
    s2 = wb.create_sheet("Summary")
    s2.append(["Metric", "Value"])
    s2.append(["Total Rows", ws.max_row - 1])
    s2.append(["Generated", datetime.utcnow().isoformat()])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _docx(name: str, r: random.Random) -> bytes:
    if docx is None:
        return _txt_fallback(name, r)
    doc = docx.Document()
    doc.add_heading(name, level=0)
    doc.add_paragraph(
        "RetailCo — Confidential. This document is part of the enterprise "
        "reference drive used for demonstration purposes."
    )
    for _ in range(r.randint(4, 9)):
        doc.add_heading(r.choice([
            "Executive Summary", "Objectives", "Scope", "Timeline",
            "Risks & Mitigations", "Budget", "Next Steps", "Appendix",
        ]), level=1)
        for _ in range(r.randint(2, 4)):
            doc.add_paragraph(_lorem(r, r.randint(30, 70)))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx(name: str, r: random.Random) -> bytes:
    if Presentation is None:
        return _txt_fallback(name, r)
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = name
    title_slide.placeholders[1].text = "RetailCo — Internal Deck"
    for _ in range(r.randint(4, 8)):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = r.choice([
            "Market Overview", "Q4 Performance", "Category Strategy",
            "Store Rollout", "Customer Insights", "Roadmap", "Financials",
        ])
        body = s.placeholders[1].text_frame
        body.text = _lorem(r, 12)
        for _ in range(r.randint(2, 4)):
            p = body.add_paragraph()
            p.text = "• " + _lorem(r, r.randint(6, 12))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf(name: str, r: random.Random) -> bytes:
    if _pdfcanvas is None:
        return _txt_fallback(name, r)
    buf = io.BytesIO()
    c = _pdfcanvas.Canvas(buf, pagesize=LETTER)
    width, height = LETTER
    for page in range(r.randint(1, 4)):
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, height - 72, name if page == 0 else f"{name} (cont.)")
        c.setFont("Helvetica", 10)
        y = height - 110
        for _ in range(r.randint(20, 34)):
            c.drawString(72, y, _lorem(r, r.randint(8, 16)))
            y -= 16
            if y < 72:
                break
        c.showPage()
    c.save()
    return buf.getvalue()


def _csv(name: str, r: random.Random) -> bytes:
    rows = ["date,region,store,product,units,revenue"]
    start = datetime(2025, 1, 1)
    for _ in range(r.randint(50, 400)):
        d = start + timedelta(days=r.randint(0, 480))
        units = r.randint(1, 900)
        rows.append(
            f"{d.strftime('%Y-%m-%d')},{r.choice(_REGIONS)},{r.choice(_STORES)},"
            f"\"{r.choice(_PRODUCTS)}\",{units},{round(units*_money(r,4,240),2)}"
        )
    return ("\n".join(rows) + "\n").encode()


def _json(name: str, r: random.Random) -> bytes:
    obj = {
        "name": name,
        "generatedAt": datetime.utcnow().isoformat() + "Z",
        "records": [
            {
                "id": r.randint(10000, 99999),
                "product": r.choice(_PRODUCTS),
                "region": r.choice(_REGIONS),
                "store": r.choice(_STORES),
                "units": r.randint(1, 500),
                "revenue": _money(r, 100, 50000),
            }
            for _ in range(r.randint(10, 60))
        ],
    }
    return json.dumps(obj, indent=2).encode()


def _md(name: str, r: random.Random) -> bytes:
    lines = [f"# {name}", "", "> RetailCo internal reference document.", ""]
    for _ in range(r.randint(3, 6)):
        lines.append("## " + r.choice([
            "Overview", "Details", "Process", "Owners", "SLAs", "Checklist",
        ]))
        lines.append("")
        for _ in range(r.randint(3, 6)):
            lines.append("- " + _lorem(r, r.randint(6, 14)))
        lines.append("")
    return ("\n".join(lines)).encode()


def _txt(name: str, r: random.Random) -> bytes:
    return _lorem(r, r.randint(80, 240)).encode()


def _txt_fallback(name: str, r: random.Random) -> bytes:
    return (f"{name}\n\n" + _lorem(r, 120)).encode()


def _png(name: str, r: random.Random) -> bytes:
    if Image is None:
        return _tiny_png()
    w, h = 640, 360
    img = Image.new("RGB", (w, h), (r.randint(20, 60), r.randint(40, 90), r.randint(80, 160)))
    d = ImageDraw.Draw(img)
    for _ in range(r.randint(8, 20)):
        x0, y0 = r.randint(0, w), r.randint(0, h)
        x1, y1 = x0 + r.randint(20, 180), y0 + r.randint(20, 120)
        d.rectangle([x0, y0, x1, y1], fill=(r.randint(0, 255), r.randint(0, 255), r.randint(0, 255)))
    d.text((16, 16), name[:40], fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _jpg(name: str, r: random.Random) -> bytes:
    if Image is None:
        return _tiny_png()
    w, h = 800, 600
    img = Image.new("RGB", (w, h), (r.randint(60, 200), r.randint(60, 200), r.randint(60, 200)))
    d = ImageDraw.Draw(img)
    for _ in range(r.randint(10, 25)):
        cx, cy = r.randint(0, w), r.randint(0, h)
        rad = r.randint(10, 90)
        d.ellipse([cx - rad, cy - rad, cx + rad, cy + rad],
                  fill=(r.randint(0, 255), r.randint(0, 255), r.randint(0, 255)))
    d.text((16, 16), name[:40], fill=(0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=82)
    return buf.getvalue()


def _tiny_png() -> bytes:
    return bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
        "0000000d49444154789c6360000002000100" "05fe02fea7d0e3450000000049454e44ae426082"
    )


_LOREM = (
    "revenue margin inventory assortment planogram markdown replenishment "
    "supplier logistics fulfillment omnichannel loyalty conversion basket "
    "shrinkage forecast promotion category vendor compliance staffing footfall "
    "clearance seasonal warehouse distribution merchandising procurement audit"
).split()


def _lorem(r: random.Random, n: int) -> str:
    words = [r.choice(_LOREM) for _ in range(n)]
    words[0] = words[0].capitalize()
    return " ".join(words) + "."


# A tiny (2s, 64x48) H.264/MP4 test-pattern clip, embedded so video files are
# genuinely playable without requiring an encoder at generation time.
_TINY_MP4_B64 = (
    "AAAAIGZ0eXBpc29tAAACAGlzb21pc28yYXZjMW1wNDEAAANMbW9vdgAAAGxtdmhkAAAAAAAAAAAA"
    "AAAAAAAD6AAAB9AAAQAAAQAAAAAAAAAAAAAAAAEAAAAAAAAAAAAAAAAAAAABAAAAAAAAAAAAAAAA"
    "AABAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAgAAAnZ0cmFrAAAAXHRraGQAAAADAAAA"
    "AAAAAAAAAAABAAAAAAAAB9AAAAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAAAAAAAAAABAAAAAAAA"
    "AAAAAAAAAABAAAAAAEAAAAAwAAAAAAAkZWR0cwAAABxlbHN0AAAAAAAAAAEAAAfQAAAAAAABAAAA"
    "AAHubWRpYQAAACBtZGhkAAAAAAAAAAAAAAAAAAAoAAAAUABVxAAAAAAALWhkbHIAAAAAAAAAAHZp"
    "ZGUAAAAAAAAAAAAAAABWaWRlb0hhbmRsZXIAAAABmW1pbmYAAAAUdm1oZAAAAAEAAAAAAAAAAAAA"
    "ACRkaW5mAAAAHGRyZWYAAAAAAAAAAQAAAAx1cmwgAAAAAQAAAVlzdGJsAAAAuXN0c2QAAAAAAAAA"
    "AQAAAKlhdmMxAAAAAAAAAAEAAAAAAAAAAAAAAAAAAAAAAEAAMABIAAAASAAAAAAAAAABAAAAAAAA"
    "AAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAGP//AAAAL2F2Y0MBQsAK/+EAF2dCwArZBHsBEAAA"
    "AwAQAAADAKDxImSAAQAFaMuA5LIAAAAQcGFzcAAAAAEAAAABAAAAFGJ0cnQAAAAAAAAV1AAAFdQA"
    "AAAYc3R0cwAAAAAAAAABAAAACgAACAAAAAAUc3RzcwAAAAAAAAABAAAAAQAAABxzdHNjAAAAAAAA"
    "AAEAAAABAAAACgAAAAEAAAA8c3RzegAAAAAAAAAAAAAACgAAA8UAAAAdAAAAMgAAADIAAAA8AAAA"
    "MgAAADQAAAA4AAAAMAAAACUAAAAUc3RjbwAAAAAAAAABAAADfAAAAGJ1ZHRhAAAAWm1ldGEAAAAA"
    "AAAAIWhkbHIAAAAAAAAAAG1kaXJhcHBsAAAAAAAAAAAAAAAALWlsc3QAAAAlqXRvbwAAAB1kYXRh"
    "AAAAAQAAAABMYXZmNTguNzYuMTAwAAAACGZyZWUAAAV9bWRhdAAAAnAGBf//bNxF6b3m2Ui3lizY"
    "INkj7u94MjY0IC0gY29yZSAxNjMgcjMwNjAgNWRiNmFhNiAtIEguMjY0L01QRUctNCBBVkMgY29k"
    "ZWMgLSBDb3B5bGVmdCAyMDAzLTIwMjEgLSBodHRwOi8vd3d3LnZpZGVvbGFuLm9yZy94MjY0Lmh0"
    "bWwgLSBvcHRpb25zOiBjYWJhYz0wIHJlZj0zIGRlYmxvY2s9MTowOjAgYW5hbHlzZT0weDE6MHgx"
    "MTEgbWU9aGV4IHN1Ym1lPTcgcHN5PTEgcHN5X3JkPTEuMDA6MC4wMCBtaXhlZF9yZWY9MSBtZV9y"
    "YW5nZT0xNiBjaHJvbWFfbWU9MSB0cmVsbGlzPTEgOHg4ZGN0PTAgY3FtPTAgZGVhZHpvbmU9MjEs"
    "MTEgZmFzdF9wc2tpcD0xIGNocm9tYV9xcF9vZmZzZXQ9LTIgdGhyZWFkcz0xIGxvb2thaGVhZF90"
    "aHJlYWRzPTEgc2xpY2VkX3RocmVhZHM9MCBucj0wIGRlY2ltYXRlPTEgaW50ZXJsYWNlZD0wIGJs"
    "dXJheV9jb21wYXQ9MCBjb25zdHJhaW5lZF9pbnRyYT0wIGJmcmFtZXM9MCB3ZWlnaHRwPTAga2V5"
    "aW50PTI1MCBrZXlpbnRfbWluPTUgc2NlbmVjdXQ9NDAgaW50cmFfcmVmcmVzaD0wIHJjX2xvb2th"
    "aGVhZD00MCByYz1jcmYgbWJ0cmVlPTEgY3JmPTQwLjAgcWNvbXA9MC42MCBxcG1pbj0wIHFwbWF4"
    "PTY5IHFwc3RlcD00IGlwX3JhdGlvPTEuNDAgYXE9MToxLjAwAIAAAAFNZYiE/+ReUsbFAB8UMN+z"
    "wBalMMiTBE+CAARs2EguNeAKluGYxjqBQ1VKoDEpYABIABfz//wyGAaGAATAgFTO//7B/37/+/wM"
    "OIwcAAkAAIC3BgEAAQKYMAAaL88IJpTHEBnzyxxQ+ARn6q7r0c2CkwqYACCABHQBQRF/X8AMw5lM"
    "DN+ug/7Qh/8EWGCvs9/n3FR5mT5T/p/9h8oj3/D92Okf8IHkhyL9iQ+eqrfIjLXl9YOBDgQFzgYF"
    "zUh0EZcvwon9AqDKQqAyksMMY52jE9VcRwAoBhJ69bO8FfLfDFnDCQLB0EZCwGJLEHpVD+HcAok8"
    "CZYEP57vAMAA/h+IH7y226nOZgYAAoH0MGsA+AdpDwGZZIF8cgWgsG0hQG0hcGjXHH9oBlD/qG4v"
    "L71NheDba/AsiGgAFkIATJqBSvmlCn44BMigePBNLFxbAAAAGUGaOeJ4fJqqrq6//Tz763+Cbu8d"
    "aOj/r5gAAAAuQZpUeJ8LExdRhlF9Hfpptm/z6lq2TWanf/wtxWfveTfbbTEd+GvN+k6bbY93HgAA"
    "AC5BmmFxPhbqpqTziNUXyd+3pptt/OQf35aX+Fu7nx+8jc3p2200/njVEFNM+v34AAAAOEGagXE+"
    "HuT9xWN+v6O/MZTfr9+GicVjKOIhMirTT/4X6i6rwGUeqbK7bZNj+F/F4W+47n451HHAAAAALkGa"
    "oJxPhbq6vv9NNsf6X/5ATk3do+dvwsXEvVKyvljOTx5fr1C/VVX/Fsxmuv4AAAAwQZrAnE+eaC2j"
    "BR6Vmpu3+Fybu75O1tndjLQQmYW/+f75TI1JGgj+F+qquOy/xbluAAAANEGa4NxPhbu7vgnP9Ofc"
    "v8K1Ifb8/Zx/DPVRjJTMZbOatnqTf4X6pjnv86EzJee6/qmQZfsAAAAsQZsARxPn3xVxb+PcdGpA"
    "sAE/gOwFaSvhog7hl74bqv7Tn3/4Z6i5/81U28AAAAAhQZsgZxPn08/bWfZv/RZvzkGqIJOmb5vJ"
    "/wU9V1dV9GOA"
)


def _mp4(name: str, r: random.Random) -> bytes:
    return base64.b64decode(_TINY_MP4_B64)


def _mp3(name: str, r: random.Random) -> bytes:
    # Valid silent MPEG-1 Layer III stream (44.1 kHz, 128 kbps): each 417-byte
    # frame is a bare frame header followed by zeroed payload (~26 ms each).
    frame = bytes([0xFF, 0xFB, 0x90, 0x64]) + bytes(413)
    seconds = r.randint(2, 5)
    return frame * (38 * seconds)


_BUILDERS = {
    "xlsx": _xlsx, "docx": _docx, "pptx": _pptx, "pdf": _pdf, "csv": _csv,
    "json": _json, "md": _md, "txt": _txt, "png": _png, "jpg": _jpg,
    "mp4": _mp4, "mp3": _mp3,
}


def build(ext: str, name: str, seed: int) -> tuple[bytes, str]:
    """Return (bytes, mime_type) for a file of type ``ext`` named ``name``."""
    ext = ext.lower()
    r = _rng(name, seed)
    builder = _BUILDERS.get(ext, _txt)
    data = builder(name, r)
    return data, MIME.get(ext, "application/octet-stream")
